"""Generate the Edge AIOps competition submission PPTX.

Run:
    python tools/generate_pptx.py

Output:
    docs/edge-aiops-submission.pptx

This script reads the structured outline in docs/ppt-outline.md and produces
a 12-slide PowerPoint deck with speaker notes.
"""
from __future__ import annotations

from pathlib import Path
from typing import List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


# ----------------------------------------------------------------------
# Color palette (from ppt-outline.md)
# ----------------------------------------------------------------------

DEEP_BLUE = RGBColor(0x0B, 0x2D, 0x5C)
MID_BLUE = RGBColor(0x1F, 0x4E, 0x8C)
ACCENT_GREEN = RGBColor(0x1E, 0x8E, 0x3E)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF8)
DARK_TEXT = RGBColor(0x20, 0x24, 0x2C)
SOFT_WHITE = RGBColor(0xFA, 0xFB, 0xFC)
TABLE_HEADER_BLUE = RGBColor(0xDD, 0xE7, 0xF2)
DANGER_RED = RGBColor(0xC0, 0x39, 0x2B)


# ----------------------------------------------------------------------
# Slide content (mirrors docs/ppt-outline.md)
# ----------------------------------------------------------------------

SLIDES: List[dict] = [
    {
        "title": "Edge AIOps",
        "subtitle": "Multi-Agent Incident Response for Edge + AI Fleets",
        "tag": "Agent Infra 初赛方案 · 2026",
        "bg": DEEP_BLUE,
        "fg": SOFT_WHITE,
        "notes": (
            "我们做的是「边缘 + 模型 + OTA」场景下的多 Agent 协同运维，名字叫 Edge AIOps。"
            "下面我会讲清楚为什么、怎么做、能落地到什么程度。"
        ),
    },
    {
        "title": "痛点：边缘 AI 部署的 4 个现实难题",
        "bullets": [
            "① 边缘节点数量大：千级 GPU 节点，间歇性断网",
            "② 模型版本多：数十个版本通过 OTA 滚动",
            "③ 漂移信号多：accuracy / concept / data / upstream 4 类",
            "④ 风险分级模糊：哪个动作能自动做？哪个必须人批？",
        ],
        "conclusion": "→ 单 Agent 串行失去阶段边界；多工具独立调用失去证据链",
        "bg": LIGHT_GREY,
        "fg": DARK_TEXT,
        "notes": (
            "边缘 AI 部署的规模已经到了千节点级别。一个模型版本通过 OTA 推到几十个区域，"
            "每个区域又有自己的客户分级。出了问题时，到底是模型回归、数据漂移、节点故障还是网络问题？"
            "一个 Agent 串行处理会失去阶段边界；多个工具独立调用会失去证据链。"
            "我们需要的是一个有清晰阶段、有风险分级、能 HITL 的多 Agent 团队。"
        ),
    },
    {
        "title": "我们的方案：4 Agent 协同团队",
        "bullets": [
            "edge_sentinel（边缘哨兵·检测）—— 融合边缘 + 模型 + OTA 三类信号",
            "edge_diagnostician（边缘诊断师·归因）—— 多源证据关联 + 轻量 RAG",
            "edge_orchestrator（边缘编排师·规划）—— 风险分级 + 发起执行请求",
            "edge_sentry（边缘验收师·验证）—— SLA 验收 + postmortem",
        ],
        "tagline": "Manager-TeamLeader-Worker 模式严格对齐 · 4 业务 Agent + 1 TeamLeader + 1 Manager = 6 Worker",
        "bg": SOFT_WHITE,
        "fg": DARK_TEXT,
        "notes": (
            "4 个 Agent，每个一个阶段。EdgeSentinel 看的是边缘 + 模型 + OTA 三类信号，做信号融合；"
            "EdgeDiagnostician 看日志、Trace、模型遥测、OTA 历史、节点资源、Runbook，做多源归因；"
            "EdgeOrchestrator 拿到归因后做风险分级，发起执行请求；"
            "EdgeSentry 验收恢复结果，写 postmortem。"
            "4 个 Agent 加 1 个 TeamLeader 加 1 个 Manager，共 6 个 Worker。和官方的 Manager-TeamLeader-Worker 模式严格对齐。"
        ),
    },
    {
        "title": "关键设计：Agents Decide, Operators Execute",
        "flow": [
            ("AgentTeams (4 Agents)", "decision", MID_BLUE),
            ("Request Queue", "hitl", ACCENT_ORANGE),
            ("External Operator Pool", "execution", ACCENT_GREEN),
        ],
        "rationale": [
            "Agent 永远不直接执行高风险动作",
            "L0/L1 → auto_approved · L2/L3 → pending_approval",
            "Demo → 生产：Agent 代码零改动，只换 mock 实现",
        ],
        "bg": SOFT_WHITE,
        "fg": DARK_TEXT,
        "notes": (
            "关键设计：Agent 不直接执行高风险动作。Agent 只发「执行请求」，带风险等级。"
            "L0 L1 自动批准，L2 L3 进入 pending_approval 状态，等人类审批。"
            "这个队列就是 HITL seam，也是 demo 跨越到生产的稳定接口。"
            "Demo 阶段我们 mock 这个队列；生产阶段由外部 Operator（K8s Operator / ML 平台 Operator / OTA Operator）监听并执行。"
            "Agent 代码永远不变。"
        ),
    },
    {
        "title": "工作流：DAG-Ready 的 dependency_graph",
        "dag": [
            ("edge_sentinel", "detect", []),
            ("edge_diagnostician", "attribute", ["detect"]),
            ("edge_orchestrator", "plan", ["attribute"]),
            ("edge_sentry", "verify", ["plan"]),
        ],
        "highlights": [
            "JSON 格式 + K8s CRD YAML 镜像同步",
            "Manager 动态调度，加新 Agent 仅 3 行 JSON 改动",
            "AgentTeams v1.2 原生支持",
        ],
        "bg": SOFT_WHITE,
        "fg": DARK_TEXT,
        "notes": (
            "工作流以 dependency_graph 形式声明，AgentTeams Manager 动态调度，"
            "加新 Agent 就是 3 行 JSON 改动。YAML 格式（K8s CRD）也同步提供，AgentTeams v1.2 原生支持。"
        ),
    },
    {
        "title": "Skill 工程体系：7 个版本化 Skill",
        "skills_table": [
            ("alert-fusion", "告警融合", "v0.1.0"),
            ("impact-mapping", "影响映射", "v0.1.0"),
            ("log-trace-rca", "多源归因", "v0.1.0"),
            ("model-governance", "模型治理", "v0.1.0"),
            ("remediation-plan", "修复计划生成", "v0.1.0"),
            ("risk-guard", "L0-L3 风险分级", "v0.1.0"),
            ("recovery-verify", "恢复验收", "v0.1.0"),
        ],
        "highlight": "每个 Skill 含：YAML frontmatter (version + maturity) · Inputs · Procedure · Output Contract · Quality Gates",
        "bg": LIGHT_GREY,
        "fg": DARK_TEXT,
        "notes": (
            "7 个 Skill，每个都是「决策类」可复用能力。结构标准化：YAML frontmatter 含版本和成熟度；"
            "输入输出有显式契约；最后有 Quality Gates 作为硬约束。"
            "比如 risk-guard 永远不自动批准 L2/L3；recovery-verify 永远要求每个 postmortem 必须含流程改进和可观测改进。"
            "这些 Quality Gates 是 Skill 工程最值钱的部分。"
        ),
    },
    {
        "title": "Scenario 1: model_drift 跑通演示",
        "trigger": "14:00 OTA 把 resnet50 v3.2.1 推到 Hangzhou 集群（无 canary）",
        "symptom": "accuracy 0.92 → 0.71, p99 420ms → 1450ms",
        "flow": [
            "① Sentinel: 3 个告警融合为 INC-2001 (P1)",
            "② Diagnostician: log-trace-rca + 轻量 RAG（score=6.0）→ top_candidate: ota_model_regression, confidence=0.86",
            "③ Orchestrator: quarantine_node L1 auto_approved · rollback_model L2 pending_approval",
            "④ Sentry: query_metrics(after) → accuracy=0.91, fleet_healthy · postmortem 自动生成",
        ],
        "evidence": "curl http://127.0.0.1:18090/tools/model_drift/requests",
        "bg": SOFT_WHITE,
        "fg": DARK_TEXT,
        "notes": (
            "我们准备了两个场景，跑通了 demo 的是模型漂移。真实还原了一个没做 canary 的 OTA 推送，"
            "14 点准时触发，准确率从 0.92 跌到 0.71。"
            "EdgeSentinel 把 3 个告警融合成一个 P1 事故；"
            "EdgeDiagnostician 调 Runbook RAG 拿到 6.0 分的 RB-OTA-001，归因到 ota_model_regression 置信度 0.86；"
            "EdgeOrchestrator 发出 2 个执行请求，L1 自动批准 L2 卡在审批；"
            "EdgeSentry 验证 L1 自动批准部分让准确率回到 0.91。"
            "整个链路可审计、可复现。"
        ),
    },
    {
        "title": "工程实现细节",
        "layout": [
            ("at/", "team_spec.json + team_spec.yaml + AgentTeam.md + RUNBOOK"),
            ("agents/", "4 个 Agent.md（含 Position + IO 契约）"),
            ("skills/", "7 个 SKILL.md（含 frontmatter + 质量门）"),
            ("tools/", "mock_tool_server.py + mock_tools.py + tool_catalog.json"),
            ("scenarios/", "2 个 JSON（model_drift + edge_node_failure）"),
            ("docs/", "submission summary + roadmap + RAG + run evidence + PPT"),
        ],
        "highlights": [
            "Mock 与真协议接口签名 1:1 对齐（tool_catalog.json 标 future_mcp_mapping）",
            "11 个 mock 工具分 3 类：decision-support / execution-request / verification",
            "任何工具的 mock → 真实现 都是 1 行改动",
            "7 个 Skill 全版本化（frontmatter 含 version + maturity）",
        ],
        "bg": LIGHT_GREY,
        "fg": DARK_TEXT,
        "notes": (
            "仓库结构清晰，4 个 Agent + 7 个 Skill + 11 个 mock 工具 + 2 个场景。"
            "重点是 mock 和真协议接口对齐，所以从 demo 跨越到生产是「替换实现」，不是「重写 Agent」。"
        ),
    },
    {
        "title": "与示例（opspilot-zero-demo）的差异",
        "diff_table": [
            ("alert-intake", "edge_sentinel", "Edge 信号"),
            ("rca-analyst", "edge_diagnostician", "节点/模型/OTA 多源归因"),
            ("remediation-planner", "edge_orchestrator", "发执行请求，不直接执行"),
            ("recovery-verifier", "edge_sentry", "SLA + 队列状态验证"),
        ],
        "innovations": [
            "命名体现 Edge + 模型 + OTA 业务",
            "Agent 不直接执行高风险动作，只发请求",
            "加 dependency_graph + K8s CRD 镜像",
            "mock_runbook 升级为轻量 RAG（tag + symptom 评分）",
            "删除对 K3s/Prometheus 的过早抽象",
        ],
        "bg": SOFT_WHITE,
        "fg": DARK_TEXT,
        "notes": (
            "我们的设计与示例有 4 点显著差异。命名反映 Edge 业务，Agent 不直接执行，依赖图显式声明，"
            "mock 实现含轻量 RAG。我们没有抄袭套壳，而是吸收示例的 Manager-Worker 模式后做了适配。"
        ),
    },
    {
        "title": "工程验证 + 可观测",
        "mock_checks": [
            "✓ Health / Scenarios / Requests / Trace 4 个端点",
            "✓ L0/L1 auto_approved → after metrics 显示恢复",
            "✓ L2/L3 pending_approval → 不执行，留在队列",
            "✓ Lightweight RAG: query score=6.0 + breakdown",
        ],
        "real_path": [
            "K3s 单机 + Triton + Prometheus + MLflow + Argo Rollouts",
            "External Operator Pool（5 ops）",
            "LoongSuite OTel 自动埋点",
            "详见 docs/final-roadmap.md（4 周时间表）",
        ],
        "observability": [
            "GET /tools/{scenario}/trace   每次工具调用",
            "GET /tools/{scenario}/requests  每个执行请求状态",
            "每个 Agent 输出含 evidence_refs",
        ],
        "bg": LIGHT_GREY,
        "fg": DARK_TEXT,
        "notes": (
            "Mock 阶段我们已经跑通了 4 类验证：健康检查、L0/L1 自动批准、L2/L3 挂起、轻量 RAG 检索。"
            "复赛的真实落地路径见 final-roadmap.md，包括 K3s 单节点、Triton 推理、Prometheus 抓取、MLflow 模型注册、Argo Rollouts OTA。"
            "可观测性靠 trace 端点和 requests 端点，每次工具调用和每个执行请求都有审计。"
        ),
    },
    {
        "title": "开源贡献计划：5 个 PR 目标",
        "prs": [
            ("agentscope-ai/AgentTeams", "docs/examples/edge-aiops.md"),
            ("agentscope-ai/AgentTeams", "Skill frontmatter 模板生成器"),
            ("agentscope-ai/MCP", "edge_aiops_mcp server（包装 11 个 mock 工具）"),
            ("agentscope-ai/AgentScope", "LoongSuite OTel Agent 埋点 hook"),
            ("argoproj/argo-rollouts", "边缘命名空间 RBAC 示例"),
        ],
        "rubric": [
            ("场景价值", "25%", "真实 Edge AIOps 场景，命中方向 1「零人工运维」"),
            ("多 Agent 协作", "25%", "4 Agent + dependency_graph"),
            ("Skill 工程", "25%", "7 Skill + frontmatter + 质量门"),
            ("工程验证", "20%", "mock 跑通 + final-roadmap"),
            ("开源贡献", "5%", "5 个 PR"),
        ],
        "bg": SOFT_WHITE,
        "fg": DARK_TEXT,
        "notes": (
            "我们准备贡献 5 个 PR，覆盖文档、Skill 模板、MCP server、可观测 hook、K8s RBAC 示例。"
            "完全对齐官方 5 个评分维度。"
        ),
    },
    {
        "title": "Q & A",
        "qas": [
            ("为什么 4 个 Agent 不多不少？",
             "严格对应「检测-归因-规划-验证」4 阶段；多一个就模糊焦点；少一个就丢阶段边界。"),
            ("为什么 Agent 不直接执行高风险动作？",
             "LLM 不可靠 + L2/L3 不可逆；通过 HITL seam 让 Agent 做参谋、外部 Operator 做执行。"),
            ("Skill 如何复用？",
             "7 个 Skill 全部决策类、命名按「决策动词」、含 frontmatter 版本化；可复用到其他 AIOps 场景。"),
            ("如何落地到生产？",
             "final-roadmap.md 详细写了 K3s + Triton + Prometheus + MLflow + Argo Rollouts + 5 个 Operator 的对接计划。"),
            ("真实数据从哪来？",
             "决赛对接真实 Triton 推理服务；可手动 push v3.2.1 复现 OTA 回归。"),
            ("与 opspilot-zero-demo 区别？",
             "命名体现 Edge 业务 · Agent 不直接执行 · dependency_graph 显式 · 轻量 RAG · 删除过早抽象。"),
        ],
        "bg": DEEP_BLUE,
        "fg": SOFT_WHITE,
        "notes": (
            "欢迎评审老师提问。我们准备的所有问题都在 PPT 里，"
            "仓库 README + docs/* 也都有详细说明。"
        ),
    },
]


# ----------------------------------------------------------------------
# Rendering helpers
# ----------------------------------------------------------------------

def add_title(slide, text: str, font_color, size: int = 32) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = font_color


def add_bullets(slide, items: List[str], color, top: float = 1.5) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.5), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_conclusion(slide, text: str, color) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(11.5), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = color


def add_tagline(slide, text: str, color) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.2), Inches(11.5), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = color


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def fill_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_table(slide, headers: List[str], rows: List[Tuple], top: float = 1.7, left: float = 0.7,
              width: float = 11.5, height: float = 4.5, header_bg=TABLE_HEADER_BLUE, header_fg=DARK_TEXT) -> None:
    rows_n = len(rows) + 1
    cols_n = len(headers)
    table_shape = slide.shapes.add_table(rows_n, cols_n, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = header_fg
        cell.text_frame.paragraphs[0].font.size = Pt(18)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].font.color.rgb = DARK_TEXT


def add_flow_diagram(slide, flow: List[Tuple[str, str, RGBColor]]) -> None:
    box_w = 3.3
    gap = 0.5
    total = len(flow)
    start_x = (13.0 - (total * box_w + (total - 1) * gap)) / 2.0
    for i, (text, _, color) in enumerate(flow):
        left = Inches(start_x + i * (box_w + gap))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(box_w), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.text = text
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = SOFT_WHITE
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = 2  # center
        if i < total - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           Inches(start_x + i * (box_w + gap) + box_w + 0.05),
                                           Inches(2.35), Inches(0.4), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_TEXT


def add_dag(slide, dag: List[Tuple[str, str, List[str]]]) -> None:
    box_w = 2.7
    gap = 0.4
    positions = [(1.0, "edge_sentinel", ACCENT_GREEN),
                 (4.1, "edge_diagnostician", MID_BLUE),
                 (7.2, "edge_orchestrator", ACCENT_ORANGE),
                 (10.3, "edge_sentry", DANGER_RED)]
    for left, name, color in positions:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.3), Inches(box_w), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.text = name
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = SOFT_WHITE
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = 2  # center

    # arrows between consecutive boxes
    for i in range(len(positions) - 1):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(positions[i][0] + box_w - 0.1),
                                       Inches(2.7), Inches(0.5), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_TEXT


# ----------------------------------------------------------------------
# Slide builders
# ----------------------------------------------------------------------

def build_slide(prs, idx: int, spec: dict) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide, spec["bg"])

    title = spec["title"]
    fg = spec["fg"]
    add_title(slide, title, fg)

    notes = spec.get("notes", "")
    add_notes(slide, notes)

    if idx == 0:
        # Cover
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.5), Inches(1.5))
        sub_p = sub_box.text_frame.paragraphs[0]
        sub_p.text = spec["subtitle"]
        sub_p.font.size = Pt(36)
        sub_p.font.bold = True
        sub_p.font.color.rgb = SOFT_WHITE

        tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.5), Inches(0.6))
        tag_p = tag_box.text_frame.paragraphs[0]
        tag_p.text = spec["tag"]
        tag_p.font.size = Pt(20)
        tag_p.font.color.rgb = SOFT_WHITE

        # decorative accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(2.0), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_GREEN

    elif idx == 1:
        add_bullets(slide, spec["bullets"], spec["fg"], top=1.6)
        add_conclusion(slide, spec["conclusion"], ACCENT_ORANGE)

    elif idx == 2:
        add_bullets(slide, spec["bullets"], spec["fg"], top=1.6)
        add_tagline(slide, spec["tagline"], ACCENT_GREEN)

    elif idx == 3:
        add_flow_diagram(slide, spec["flow"])
        # rationale bullets below
        add_bullets(slide, spec["rationale"], spec["fg"], top=4.0)

    elif idx == 4:
        add_dag(slide, spec["dag"])
        add_bullets(slide, spec["highlights"], spec["fg"], top=4.0)

    elif idx == 5:
        add_table(slide, ["Skill 名", "职责", "版本"], spec["skills_table"], top=1.6, height=4.2)
        # highlight below table
        hl_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(11.5), Inches(0.7))
        hl_p = hl_box.text_frame.paragraphs[0]
        hl_p.text = spec["highlight"]
        hl_p.font.size = Pt(16)
        hl_p.font.italic = True
        hl_p.font.color.rgb = ACCENT_GREEN

    elif idx == 6:
        # trigger + symptom boxes
        for i, (label, value) in enumerate([("Trigger:", spec["trigger"]), ("Symptom:", spec["symptom"])]):
            box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4 + i * 0.6), Inches(11.5), Inches(0.5))
            p = box.text_frame.paragraphs[0]
            p.text = f"{label} {value}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            if i == 0:
                p.font.bold = True
        add_bullets(slide, spec["flow"], DARK_TEXT, top=3.0)
        # evidence
        ev_box = slide.shapes.add_textbox(Inches(0.7), Inches(7.3), Inches(11.5), Inches(0.4))
        ev_p = ev_box.text_frame.paragraphs[0]
        ev_p.text = f"复现命令: {spec['evidence']}"
        ev_p.font.size = Pt(14)
        ev_p.font.italic = True
        ev_p.font.color.rgb = ACCENT_GREEN

    elif idx == 7:
        add_table(slide, ["目录", "内容"], spec["layout"], top=1.6, height=3.0)
        # highlights bullets
        add_bullets(slide, spec["highlights"], DARK_TEXT, top=5.2)

    elif idx == 8:
        add_table(slide, ["opspilot-zero-demo", "edge-aiops-demo (我们的)", "差异"], spec["diff_table"], top=1.6, height=2.5)
        add_bullets(slide, spec["innovations"], DARK_TEXT, top=4.5)

    elif idx == 9:
        # two-column layout: mock_checks / real_path
        box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5))
        tf1 = box1.text_frame
        tf1.text = "Mock 验证（已跑通）"
        tf1.paragraphs[0].font.size = Pt(20)
        tf1.paragraphs[0].font.bold = True
        tf1.paragraphs[0].font.color.rgb = ACCENT_GREEN
        for item in spec["mock_checks"]:
            p = tf1.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(4)

        box2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5))
        tf2 = box2.text_frame
        tf2.text = "真实落地路径（复赛）"
        tf2.paragraphs[0].font.size = Pt(20)
        tf2.paragraphs[0].font.bold = True
        tf2.paragraphs[0].font.color.rgb = ACCENT_ORANGE
        for item in spec["real_path"]:
            p = tf2.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(4)

        # observability row
        obs_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.5))
        obs_p = obs_box.text_frame.paragraphs[0]
        obs_p.text = "可观测:  " + "  ·  ".join(spec["observability"])
        obs_p.font.size = Pt(14)
        obs_p.font.italic = True
        obs_p.font.color.rgb = DARK_TEXT

    elif idx == 10:
        add_table(slide, ["仓库", "PR 内容"], spec["prs"], top=1.6, height=3.5)
        # rubric below
        add_bullets(slide, [f"{w}  ({pct}):  {desc}" for w, pct, desc in spec["rubric"]],
                    DARK_TEXT, top=5.4)

    elif idx == 11:
        # Q&A
        add_table(slide, ["问题", "回答"], spec["qas"], top=1.6, height=5.5, header_bg=DEEP_BLUE, header_fg=SOFT_WHITE)


# ----------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, spec in enumerate(SLIDES):
        build_slide(prs, idx, spec)

    output = Path(__file__).resolve().parents[1] / "docs" / "edge-aiops-submission.pptx"
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Saved: {output}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()